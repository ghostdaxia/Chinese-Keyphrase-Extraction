import PyPDF2
import docx2txt
import pandas as pd
from pptx import Presentation

def read_pdf_to_text(file_path):
    with open(file_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        contents_list = []
        for page in pdf_reader.pages:
            content = page.extract_text()
            contents_list.append(content)
    return '\n'.join(contents_list)


def read_docx_to_text(file_path):
    text = docx2txt.process(file_path)
    return text


def read_excel_to_text(file_path):
    excel_file = pd.ExcelFile(file_path)
    sheet_names = excel_file.sheet_names

    text_list = []
    for sheet_name in sheet_names:
        df = excel_file.parse(sheet_name)
        text = df.to_string(index=False)
        text_list.append(text)
    return '\n'.join(text_list)


def read_pptx_to_text(file_path):
    prs = Presentation(file_path)

    text_list = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text = text_frame.text
                if text:
                    text_list.append(text)
    return '\n'.join(text_list)

def read_txt_to_text(file_path):
    with open(file_path, 'r') as f:
        text = f.read()
    return text


support = {
    'pdf': 'read_pdf_to_text',
    'docx': 'read_docx_to_text',
    'xlsx': 'read_excel_to_text',
    'xls': 'read_excel_to_text',
    'pptx': 'read_pptx_to_text',
    'ppt': 'read_pptx_to_text',
    'csv': 'read_txt_to_text',
    'txt': 'read_txt_to_text'
}

def read_any_file_to_text(file_path):
    file_suffix = file_path.split('.')[-1]
    func = support.get(file_suffix)
    if func is None:
        return '暂不支持该文件格式'
    text =  eval(func)(file_path)
    return text

